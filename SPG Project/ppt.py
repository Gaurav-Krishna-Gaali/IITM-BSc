from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Case Study Project Team Roles and Responsibilities"
subtitle.text = "An Overview of Our Team’s Functions"

# Slide with condensed roles
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Team Roles and Responsibilities"

content = (
    "Project Manager:\n"
    "- Maintain project flow\n"
    "- Ensure deadlines and milestones\n"
    "- Distribute tasks and oversee progress\n\n"
    
    "Primary Research Lead:\n"
    "- Formulate questionnaires\n"
    "- Research case studies, podcasts, videos\n"
    "- Design questions for insights into EQ\n\n"
    
    "Secondary Research Lead:\n"
    "- Support Primary Lead\n"
    "- Collect and curate online resources\n"
    "- Provide new insights and references\n\n"
    
    "Data Analyst:\n"
    "- Collect and analyze data\n"
    "- Create infographics and visual aids\n"
    "- Include secondary statistics\n\n"
    
    "Writer/Editor:\n"
    "- Manage content flow\n"
    "- Maintain report quality\n"
    "- Edit infographics\n\n"
    
    "Quality Assurance:\n"
    "- Final check for quality and formatting\n"
    "- Verify infographics and citations\n"
    "- Ensure adherence to standards\n"
)

text_box = slide_2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
text_frame = text_box.text_frame
text_frame.text = content

# Save the presentation
prs.save('Case_Study_Team_Roles.pptx')
